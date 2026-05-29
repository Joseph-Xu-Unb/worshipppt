"""Generate worship PPT files from JSON data with a small Tkinter GUI."""

import json
import os
import logging
import tkinter as tk
from dataclasses import dataclass
from datetime import date
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# ================= 核心逻辑函数 =================

# 配置日志
logging.basicConfig(
    filename="worship_ppt.log",
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)


class WorshipPPTGenerator:
    """Build a worship presentation from template layouts and structured data."""

    TITLE_PLACEHOLDER_IDXS = [0, 11]
    BODY_PLACEHOLDER_IDXS = [1, 10]
    PAGE_PLACEHOLDER_IDXS = [12]

    def __init__(self, template_path, data, selected_date=None):
        """Initialize generator state and load the template presentation."""
        self.template_path = template_path
        self.data = data
        self.selected_date = selected_date
        self.prs = Presentation(template_path)
        self.warnings = []

    def _warn(self, msg):
        """Store and log a non-fatal warning message."""
        logging.warning(msg)
        self.warnings.append(msg)

    def _find_layout(self, layout_name):
        """Find a layout by name, ignoring case and surrounding spaces."""
        return next(
            (
                layout
                for layout in self.prs.slide_layouts
                if layout.name.lower().strip() == layout_name.lower().strip()
            ),
            None,
        )

    @staticmethod
    def _split_lines(lines, n=4):
        """Split text lines into chunks for paged slides."""
        return [lines[i : i + n] for i in range(0, len(lines), n)]

    @staticmethod
    def _safe_fill(slide, text, idx_list):
        """Fill the first matching placeholder index with text."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx in idx_list:
                ph.text = text
                return True
        return False

    def _add_layout_slide(self, layout_name):
        """Add one slide from a named layout, or record a warning if missing."""
        layout = self._find_layout(layout_name)
        if not layout:
            self._warn(f"Layout '{layout_name}' not found.")
            return None
        return self.prs.slides.add_slide(layout)

    def _add_paged_content(self, layout_name, title, lines):
        """Add one or more content slides with shared title and paged body text."""
        layout = self._find_layout(layout_name)
        if not layout:
            self._warn(f"Layout '{layout_name}' not found.")
            return

        for chunk in self._split_lines(lines, 4):
            slide = self.prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = title
            else:
                self._safe_fill(slide, title, self.TITLE_PLACEHOLDER_IDXS)

            content_text = "\n".join(chunk)
            self._safe_fill(slide, content_text, self.BODY_PLACEHOLDER_IDXS)
            self._safe_fill(
                slide, str(len(self.prs.slides)), self.PAGE_PLACEHOLDER_IDXS
            )

    def _rename_custom_layouts(self):
        """Rename generic custom layout names to clearer English labels."""
        layout_rename_map = {
            "8_Custom Layout": "Special Announcement",
            "9_Custom Layout": "Event Highlight",
            "13_Custom Layout": "Closing Remarks",
            "1_Custom Layout": "General Purpose",
        }
        for layout in self.prs.slide_layouts:
            if layout.name in layout_rename_map:
                layout.name = layout_rename_map[layout.name]

    def _build_order_text(self):
        """Construct the worship order text block for the order slide."""
        raw_items = [
            f"宣召经文  {self.data['call_to_worship']['title']}",
            "赞美祷告  Prayer",
        ]
        for i, hymn in enumerate(self.data["hymns"]):
            raw_items.append(f"赞美诗 {i+1}  {hymn['title']}")

        raw_items.extend(
            [
                "认罪祷告  Prayer",
                f"主题经文  {self.data['theme_scripture']['title']}",
                "讲道  Preaching",
                f"回应诗  {self.data['response_hymn']['title']}",
                "奉献  Offering",
                "三一颂  Doxology",
                "牧师祝福  Benediction",
                "事项报告  Announcements",
                "我今天为你祝福  Blessings",
            ]
        )
        return "\n".join(raw_items)

    def _add_cover(self):
        """Reuse existing template cover slide; add one only if missing."""
        cover_slide = next(
            (
                s
                for s in self.prs.slides
                if s.slide_layout.name.lower().strip() == "cover"
            ),
            None,
        )
        if not cover_slide:
            cover_slide = self._add_layout_slide("cover")

        if not cover_slide or not self.selected_date:
            return

        if cover_slide.shapes.title:
            title_text = cover_slide.shapes.title.text or ""
            if self.selected_date not in title_text:
                cover_slide.shapes.title.text = f"{title_text}\n{self.selected_date}"
        else:
            self._safe_fill(
                cover_slide, self.selected_date, self.TITLE_PLACEHOLDER_IDXS
            )

    def _add_opening_slides(self):
        """Add opening fixed slides shown before the order slide."""
        for name in ["prepare", "Slogan"]:
            self._add_layout_slide(name)

    def _add_order_slide(self):
        """Add the worship order slide with title and generated order text."""
        order_slide = self._add_layout_slide("order")
        if not order_slide:
            return
        self._safe_fill(order_slide, "敬拜次序", self.TITLE_PLACEHOLDER_IDXS)
        self._safe_fill(
            order_slide, self._build_order_text(), self.BODY_PLACEHOLDER_IDXS
        )

    def _add_call_to_worship(self):
        """Add paged call-to-worship scripture slides."""
        call_title = f"宣召经文：{self.data['call_to_worship']['title']}"
        self._add_paged_content(
            "call_to_scripture",
            call_title,
            self.data["call_to_worship"]["lines"],
        )

    def _add_hymn_section(self):
        """Add praise prayer slide and paged hymn lyric slides."""
        self._add_layout_slide("praise_prayer")
        for i, hymn in enumerate(self.data["hymns"]):
            hymn_display_title = f"Hymn{i+1}: {hymn['title']}"
            self._add_paged_content("Hymn", hymn_display_title, hymn["lines"])

    def _add_theme_scripture(self):
        """Add intercessory prayer slide and paged theme scripture slides."""
        self._add_layout_slide("intercessory_prayer")

        self._add_paged_content(
            "theme_scripture",
            f"主题经文：{self.data['theme_scripture']['title']}",
            self.data["theme_scripture"]["lines"],
        )

    def _add_sermon_slide(self):
        """Add sermon title slide."""
        s_slide = self._add_layout_slide("sermon")
        if s_slide and s_slide.shapes.title:
            s_slide.shapes.title.text = "讲道：绝境逢生"

    def _add_response_hymn(self):
        """Add paged response hymn slides."""
        response_title = f"回应诗: {self.data['response_hymn']['title']}"
        self._add_paged_content(
            "response", response_title, self.data["response_hymn"]["lines"]
        )

    def _add_static_ending(self):
        """Add fixed ending slides in sequence."""
        for name in [
            "Offering",
            "praying",
            "ode_to_the_Trinity",
            "benediction",
            "child_pickup_reminder",
            "announcements_welcome_banner",
            "matters",
            "结束",
            "wishyouwell",
        ]:
            slide = self._add_layout_slide(name)
            if not slide:
                available = [l.name for l in self.prs.slide_layouts]
                logging.warning("Available layouts in template: %s", available)
                continue

            # Provide default text for child pickup reminder page.
            if name == "child_pickup_reminder":
                title_shape = slide.shapes.title
                if not title_shape:
                    for ph in slide.placeholders:
                        if ph.placeholder_format.idx in self.TITLE_PLACEHOLDER_IDXS:
                            title_shape = ph
                            break

                if title_shape and hasattr(title_shape, "text_frame"):
                    tf = title_shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = "回归时间"
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(54)
                    run.font.bold = True

                body_shape = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 13:
                        body_shape = ph
                        break
                if not body_shape:
                    for ph in slide.placeholders:
                        if ph.placeholder_format.idx in self.BODY_PLACEHOLDER_IDXS:
                            body_shape = ph
                            break

                if body_shape and hasattr(body_shape, "text_frame"):
                    tf = body_shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = "请有小孩的弟兄姊妹和朋友到二楼领回您的孩子。"
                    p.level = 0
                    run = p.runs[0]
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(48)

            # Keep one blank announcements page after the welcome banner.
            if name == "matters":
                for ph in slide.placeholders:
                    if hasattr(ph, "text"):
                        ph.text = ""

    def generate(self, output_path):
        """Run the full generation pipeline and save the output PPT file."""
        logging.info("Start generating PPT. Template: %s", self.template_path)

        self._add_cover()
        self._add_opening_slides()
        self._add_order_slide()
        self._add_call_to_worship()
        self._add_hymn_section()
        self._add_theme_scripture()
        self._add_sermon_slide()
        self._add_response_hymn()
        self._add_static_ending()
        self._rename_custom_layouts()

        self.prs.save(output_path)
        logging.info("PPT saved to %s", output_path)
        return self.warnings


def generate_worship_ppt(template_path, data, output_path, selected_date=None):
    """Compatibility wrapper that generates PPT using the OOP generator."""
    generator = WorshipPPTGenerator(
        template_path=template_path,
        data=data,
        selected_date=selected_date,
    )
    return generator.generate(output_path)


@dataclass
class AppState:
    """Mutable state for selected files and date in the GUI workflow."""

    template_path: str | None = None
    json_path: str | None = None
    selected_date: str = ""


class ChurchAppUI:
    """Tkinter view layer: widgets and simple UI updates only."""

    def __init__(self, master, default_date):
        """Create and layout all GUI controls."""
        tk.Label(
            master,
            text="Fredericton BSBC 崇拜助手",
            font=("Arial", 18, "bold"),
            pady=20,
        ).pack()

        self.btn_tpl = tk.Button(master, text="1. 选择 PPT 模板 (.pptx)", width=25)
        self.btn_tpl.pack(pady=10)
        self.lbl_tpl = tk.Label(master, text="未选择模板", fg="grey")
        self.lbl_tpl.pack()

        self.btn_jsn = tk.Button(master, text="2. 选择 JSON 歌词数据", width=25)
        self.btn_jsn.pack(pady=10)
        self.lbl_jsn = tk.Label(master, text="未选择 JSON", fg="grey")
        self.lbl_jsn.pack()

        tk.Label(master, text="3. 选择日期 (生成文件名)", font=("Arial", 10)).pack(
            pady=5
        )
        self.date_entry = tk.Entry(master, width=15, font=("Arial", 12))
        self.date_entry.pack(pady=5)
        self.date_entry.insert(0, default_date)

        self.btn_run = tk.Button(
            master,
            text="🚀 生成完整崇拜 PPT",
            state="disabled",
            width=30,
            height=2,
            bg="#4CAF50",
            fg="white",
            font=("Arial", 11, "bold"),
        )
        self.btn_run.pack(pady=30)

    def bind_actions(self, on_tpl, on_jsn, on_run):
        """Bind button callbacks from controller methods."""
        self.btn_tpl.config(command=on_tpl)
        self.btn_jsn.config(command=on_jsn)
        self.btn_run.config(command=on_run)

    def set_template_name(self, path):
        """Show selected template filename."""
        self.lbl_tpl.config(text=os.path.basename(path), fg="black")

    def set_json_name(self, path):
        """Show selected JSON filename."""
        self.lbl_jsn.config(text=os.path.basename(path), fg="black")

    def set_run_enabled(self, enabled):
        """Enable or disable the run button."""
        self.btn_run.config(state="normal" if enabled else "disabled")

    def get_date(self):
        """Read and normalize date input from the UI."""
        return self.date_entry.get().strip()


class ChurchAppController:
    """Controller layer: orchestrates user actions and generation flow."""

    def __init__(self, master):
        """Initialize app state, root window settings, and bindings."""
        self.root = master
        self.root.title("BSBC PPT 自动化工具 v5")
        self.root.geometry("550x400")

        self.state = AppState(selected_date=date.today().isoformat())
        self.ui = ChurchAppUI(master, self.state.selected_date)
        self.ui.bind_actions(self.load_tpl, self.load_jsn, self.execute)

    def load_tpl(self):
        """Select template file and update state/UI."""
        path = filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])
        if not path:
            return
        self.state.template_path = path
        self.ui.set_template_name(path)
        self.check_ready()

    def load_jsn(self):
        """Select JSON data file and update state/UI."""
        path = filedialog.askopenfilename(filetypes=[("JSON", "*.json")])
        if not path:
            return
        self.state.json_path = path
        self.ui.set_json_name(path)
        self.check_ready()

    def check_ready(self):
        """Enable generation when both template and data file are selected."""
        ready = bool(self.state.template_path and self.state.json_path)
        self.ui.set_run_enabled(ready)

    def _load_data(self):
        """Load worship data JSON from the selected file."""
        with open(self.state.json_path, "r", encoding="utf-8") as f:
            return json.load(f)

    def _output_path(self, selected_date):
        """Build output path using JSON folder and selected date."""
        return os.path.join(
            os.path.dirname(self.state.json_path),
            f"Worship_{selected_date}.pptx",
        )

    def execute(self):
        """Validate inputs, generate PPT, and show result message."""
        try:
            selected_date = self.ui.get_date()
            if not selected_date:
                raise ValueError("日期不能为空！")

            self.state.selected_date = selected_date
            data = self._load_data()
            output_file = self._output_path(selected_date)

            warnings = generate_worship_ppt(
                self.state.template_path,
                data,
                output_file,
                selected_date=selected_date,
            )

            if warnings:
                warn_msg = "\n".join(warnings)
                messagebox.showwarning(
                    "完成 (含警告)",
                    f"PPT 已生成，但有以下问题：\n{warn_msg}\n\n保存在：\n{output_file}",
                )
            else:
                messagebox.showinfo("成功", f"PPT 已生成！\n\n保存在：\n{output_file}")
        except (ValueError, OSError, json.JSONDecodeError) as e:
            logging.error("Generation failed", exc_info=True)
            messagebox.showerror(
                "生成失败", f"错误详情：\n{str(e)}\n\n已记录到日志文件。"
            )


# Backward-compatible name
class ChurchApp(ChurchAppController):
    """Compatibility alias for existing imports and startup code."""

    pass


if __name__ == "__main__":
    app_root = tk.Tk()
    app = ChurchApp(app_root)
    app_root.mainloop()
